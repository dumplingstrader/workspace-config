"""
Create Leadership Presentation PowerPoint Template
===================================================
Generates a professional PowerPoint template for presenting
technical assistance value to leadership.

Can generate either:
1. Blank template with placeholders (default)
2. Auto-filled presentation from quarterly insights data

Usage:
    # Blank template
    python scripts/create_leadership_presentation_template.py
    
    # Auto-filled from data
    python scripts/create_leadership_presentation_template.py --quarter 2025-Q4 --input data/master_combined.json --output output/Q4_2025_Presentation.pptx

Outputs: templates/Leadership_Presentation_Template.pptx (or custom path)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import argparse
import json
import pandas as pd


# Design colors - Professional teal and gray palette
PRIMARY_COLOR = RGBColor(94, 168, 167)      # Teal #5EA8A7
SECONDARY_COLOR = RGBColor(39, 120, 132)   # Deep Teal #277884
ACCENT_COLOR = RGBColor(254, 68, 71)       # Coral #FE4447
DARK_TEXT = RGBColor(28, 40, 51)           # Dark Blue-Gray #1C2833
LIGHT_BG = RGBColor(244, 246, 246)         # Off-White #F4F6F6


def load_quarterly_data(json_path: Path, year: int, quarter: int):
    """Load and filter quarterly data from master JSON."""
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    df = pd.DataFrame(data.get('entries', []))
    df['Date'] = pd.to_datetime(df['Date'], errors='coerce')
    
    # Filter to quarter
    quarter_months = {1: [1, 2, 3], 2: [4, 5, 6], 3: [7, 8, 9], 4: [10, 11, 12]}
    months = quarter_months[quarter]
    df_quarter = df[(df['Date'].dt.year == year) & (df['Date'].dt.month.isin(months))]
    
    # Calculate metrics
    metrics = {
        'total_issues': len(df_quarter),
        'avg_per_month': len(df_quarter) / 3,
        'systems_count': df_quarter['System'].nunique() if 'System' in df_quarter.columns else 0,
        'quarter_name': f"Q{quarter} {year}",
        'year': year,
        'quarter': quarter
    }
    
    # Monthly breakdown
    df_quarter['Month'] = df_quarter['Date'].dt.to_period('M')
    monthly = df_quarter.groupby('Month').size().to_dict()
    metrics['monthly'] = {str(k): v for k, v in monthly.items()}
    
    # Top system
    if 'System' in df_quarter.columns and len(df_quarter) > 0:
        top_system = df_quarter['System'].value_counts()
        if len(top_system) > 0:
            metrics['top_system'] = top_system.index[0]
            metrics['top_system_count'] = int(top_system.iloc[0])
    
    # Complexity breakdown
    if 'Complexity' in df_quarter.columns:
        complexity = df_quarter['Complexity'].value_counts().to_dict()
        metrics['complexity'] = {k: int(v) for k, v in complexity.items()}
        metrics['quick_count'] = int(df_quarter[df_quarter['Complexity'].str.contains('Quick|<1hr', case=False, na=False)].shape[0])
    
    # Trend
    if len(metrics['monthly']) > 1:
        months_list = sorted(metrics['monthly'].items())
        first_month = months_list[0][1]
        last_month = months_list[-1][1]
        if last_month > first_month:
            metrics['trend'] = 'Increasing'
            metrics['trend_pct'] = int(((last_month - first_month) / first_month) * 100) if first_month > 0 else 0
        else:
            metrics['trend'] = 'Decreasing'
            metrics['trend_pct'] = 0
    else:
        metrics['trend'] = 'Stable'
        metrics['trend_pct'] = 0
    
    return metrics


def add_title_slide(prs, data=None):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Technical Assistance Value Report' if not data else 'Process Controls Technical Assistance'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(9), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    if data:
        subtitle_frame.text = f"{data['quarter_name']} Summary Report"
    else:
        subtitle_frame.text = '[Period: Month/Quarter YYYY]'
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Presenter name
    name_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(9), Inches(0.5)
    )
    name_frame = name_box.text_frame
    name_frame.text = '[Presenter Name] | Process Controls'
    name_para = name_frame.paragraphs[0]
    name_para.font.size = Pt(18)
    name_para.font.color.rgb = RGBColor(255, 255, 255)
    name_para.alignment = PP_ALIGN.CENTER


def add_executive_summary(prs, data=None):
    """Slide 2: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Executive Summary'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    if data:
        # Data-driven content
        trend_desc = f"{data['trend']} ({data['trend_pct']:+.0f}%)" if data['trend_pct'] != 0 else data['trend']
        top_system_text = f"{data['top_system']} ({data['top_system_count']} issues)"
        
        content = [
            ('Key Achievements', f"Handled {data['total_issues']} technical assistance requests across {data['systems_count']} systems. {data.get('quick_count', 0)} quick-response issues resolved within 24 hours."),
            ('Notable Trends', f"Workload trend: {trend_desc}. Primary system: {top_system_text}. Activity averaged {data['avg_per_month']:.1f} issues per month."),
            ('Strategic Impact', 'Maintained operational continuity across multiple sites. Reduced production impact through rapid response. Strengthened technical expertise across control platforms.')
        ]
    else:
        # Placeholder content
        content = [
            ('Key Achievements', '[Highlight top 2-3 accomplishments]'),
            ('Notable Trends', '[1-2 patterns observed]'),
            ('Strategic Impact', '[How this work supports business goals]')
        ]
    
    y_pos = 1.5
    for heading, placeholder in content:
        # Heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(y_pos),
            Inches(8.4), Inches(0.4)
        )
        heading_frame = heading_box.text_frame
        heading_frame.text = heading
        heading_para = heading_frame.paragraphs[0]
        heading_para.font.size = Pt(20)
        heading_para.font.bold = True
        heading_para.font.color.rgb = SECONDARY_COLOR
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(y_pos + 0.45),
            Inches(8.4), Inches(1.0)
        )
        content_frame = content_box.text_frame
        content_frame.text = placeholder
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(16)
        content_para.font.color.rgb = DARK_TEXT
        
        y_pos += 1.7


def add_metrics_dashboard(prs, data=None):
    """Slide 3: Key Metrics Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Key Metrics'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Metric cards (2 rows x 3 columns)
    if data:
        # Calculate additional metrics
        high_complexity = data['complexity'].get('Major', 0) + data['complexity'].get('Significant', 0)
        quick_pct = (data.get('quick_count', 0) / data['total_issues'] * 100) if data['total_issues'] > 0 else 0
        
        metrics = [
            ('Total Issues', str(data['total_issues']), 'Technical requests handled'),
            ('Systems Supported', str(data['systems_count']), 'Different control systems'),
            ('High Complexity', str(high_complexity), 'Major/significant issues'),
            ('Quick Response', f"{data.get('quick_count', 0)}", f'{quick_pct:.0f}% under 24 hours'),
            ('Monthly Average', f"{data['avg_per_month']:.1f}", 'Issues per month'),
            ('Trend', data['trend'], f"{data['trend_pct']:+.0f}% change" if data['trend_pct'] != 0 else 'Stable')
        ]
    else:
        metrics = [
            ('Total Issues', '[#]', 'Technical requests handled'),
            ('Systems Supported', '[#]', 'Different control systems'),
            ('High Complexity', '[#]', 'Major/significant issues'),
            ('Cross-Site Support', '[#]', 'Multi-site coordination'),
            ('Avg Response Time', '[X hours]', 'Time to resolution'),
            ('Uptime Protected', '[X%]', 'Production impact avoided')
        ]
    
    row = 0
    col = 0
    for i, (label, value, description) in enumerate(metrics):
        x = 0.5 + col * 3.2
        y = 1.5 + row * 2.3
        
        # Card background
        card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y),
            Inches(3.0), Inches(2.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(2)
        
        # Value (big number)
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.3),
            Inches(2.6), Inches(0.7)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = ACCENT_COLOR
        value_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.0),
            Inches(2.6), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True
        label_para.font.color.rgb = DARK_TEXT
        label_para.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.4),
            Inches(2.6), Inches(0.4)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(10)
        desc_para.font.color.rgb = DARK_TEXT
        desc_para.alignment = PP_ALIGN.CENTER
        
        col += 1
        if col == 3:
            col = 0
            row += 1


def add_system_breakdown(prs, data=None):
    """Slide 4: System Breakdown Chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Issues by System'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    if data and 'systems' in data and data['systems']:
        # Create simple bar chart representation using text
        chart_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.5),
            Inches(7), Inches(5)
        )
        chart_frame = chart_box.text_frame
        chart_frame.word_wrap = True
        
        # Sort systems by count
        sorted_systems = sorted(data['systems'].items(), key=lambda x: x[1], reverse=True)
        top_systems = sorted_systems[:10]  # Show top 10
        
        chart_text = "Top Control Systems by Issue Volume:\n\n"
        for system, count in top_systems:
            pct = (count / data['total_issues'] * 100) if data['total_issues'] > 0 else 0
            bar = "█" * int(pct / 2)  # Visual bar
            chart_text += f"{system:20s} {count:3d} issues {bar}\n"
        
        chart_para = chart_frame.paragraphs[0]
        chart_para.text = chart_text
        chart_para.font.size = Pt(14)
        chart_para.font.name = 'Consolas'
        chart_para.font.color.rgb = DARK_TEXT
    else:
        # Chart placeholder
        chart_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.0),
            Inches(7), Inches(4)
        )
        chart_frame = chart_box.text_frame
        chart_frame.text = '[INSERT CHART HERE]\n\n' \
                          'Recommended: Bar chart or pie chart showing\n' \
                          'distribution of issues across control systems'
        chart_para = chart_frame.paragraphs[0]
        chart_para.font.size = Pt(18)
        chart_para.font.italic = True
        chart_para.font.color.rgb = SECONDARY_COLOR
        chart_para.alignment = PP_ALIGN.CENTER


def add_success_stories(prs, data=None):
    """Slide 5: Success Stories / Impact Examples."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Success Stories'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Story cards
    if data:
        stories = [
            ('High Volume Response', f"Successfully handled {data['total_issues']} technical assistance requests across {data['systems_count']} different control systems during {data['quarter_name']}.",
             f"Impact: {data.get('quick_count', 0)} quick-response issues resolved within 24 hours, minimizing production disruptions."),
            ('System Expertise', f"Primary support provided for {data['top_system']} system with {data['top_system_count']} issues addressed.",
             'Impact: Demonstrated deep expertise in critical control systems, ensuring operational continuity.'),
            ('Workload Management', f"Handled {data['trend']} workload trend ({data['trend_pct']:+.0f}% change).",
             'Impact: Maintained responsive support despite changing demand, adaptable resource allocation.')
        ]
    else:
        stories = [
            ('Major Issue Resolved', '[Brief description of complex problem solved]',
             'Impact: [Production saved, downtime avoided, cost savings]'),
            ('Cross-Site Collaboration', '[Example of multi-site coordination]',
             'Impact: [Standardization achieved, best practices shared]'),
            ('Training & Knowledge Transfer', '[How you enabled others]',
             'Impact: [Team capability improved, self-sufficiency increased]')
        ]
    
    y_pos = 1.5
    for title, description, impact in stories:
        # Story card background
        card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(y_pos),
            Inches(8.4), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = SECONDARY_COLOR
        card.line.width = Pt(1)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos + 0.1),
            Inches(8.0), Inches(0.3)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = SECONDARY_COLOR
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos + 0.45),
            Inches(8.0), Inches(0.35)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = DARK_TEXT
        
        # Impact
        impact_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos + 0.85),
            Inches(8.0), Inches(0.3)
        )
        impact_frame = impact_box.text_frame
        impact_frame.text = impact
        impact_para = impact_frame.paragraphs[0]
        impact_para.font.size = Pt(12)
        impact_para.font.italic = True
        impact_para.font.color.rgb = ACCENT_COLOR
        
        y_pos += 1.5


def add_recommendations(prs, data=None):
    """Slide 6: Recommendations / Call to Action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = 'Recommendations'
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    # Recommendations list
    if data:
        recommendations = []
        
        # Trend-based recommendations
        if data['trend'] == 'Increasing':
            recommendations.append(f"Resource Planning: Workload increased {data['trend_pct']:.0f}% - consider capacity expansion")
        elif data['trend'] == 'Decreasing':
            recommendations.append(f"Process Improvement: Workload decreased {abs(data['trend_pct']):.0f}% - identify and replicate successful practices")
        
        # System-based recommendations
        if data['top_system_count'] > data['total_issues'] * 0.3:
            recommendations.append(f"Targeted Training: {data['top_system']} accounts for {data['top_system_count']} issues - invest in specialized training")
        
        # Always include standard recommendations
        recommendations.extend([
            'Continue documenting technical assistance activities for trend analysis',
            'Schedule quarterly reviews to assess workload patterns and resource needs'
        ])
        
        # Pad to 4 recommendations
        while len(recommendations) < 4:
            recommendations.append('[Additional recommendation based on specific observations]')
    else:
        recommendations = [
            'Invest in [training/equipment/resources]',
            'Standardize [process/documentation]',
            'Expand support to [new area/system]',
            'Schedule regular [reviews/maintenance]'
        ]
    
    y_pos = 1.8
    for i, rec in enumerate(recommendations, 1):
        # Number circle
        circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(1.0), Inches(y_pos),
            Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_COLOR
        circle.line.color.rgb = ACCENT_COLOR
        
        # Number text
        num_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos),
            Inches(0.4), Inches(0.4)
        )
        num_frame = num_box.text_frame
        num_frame.text = str(i)
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(18)
        num_para.font.bold = True
        num_para.font.color.rgb = RGBColor(255, 255, 255)
        num_para.alignment = PP_ALIGN.CENTER
        
        # Recommendation text
        text_box = slide.shapes.add_textbox(
            Inches(1.6), Inches(y_pos),
            Inches(7.2), Inches(0.6)
        )
        text_frame = text_box.text_frame
        text_frame.text = rec
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(18)
        text_para.font.color.rgb = DARK_TEXT
        
        y_pos += 0.9


def add_questions_slide(prs, data=None):
    """Slide 7: Questions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    # Questions text
    questions_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.0),
        Inches(9), Inches(1.5)
    )
    questions_frame = questions_box.text_frame
    questions_frame.text = 'Questions?'
    questions_para = questions_frame.paragraphs[0]
    questions_para.font.size = Pt(60)
    questions_para.font.bold = True
    questions_para.font.color.rgb = RGBColor(255, 255, 255)
    questions_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.0),
        Inches(9), Inches(0.5)
    )
    contact_frame = contact_box.text_frame
    contact_frame.text = '[Your Name] | [Email] | [Phone]'
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(18)
    contact_para.font.color.rgb = RGBColor(255, 255, 255)
    contact_para.alignment = PP_ALIGN.CENTER


def main():
    """Generate PowerPoint presentation template or auto-filled presentation."""
    
    # Parse command-line arguments
    parser = argparse.ArgumentParser(
        description='Generate PowerPoint leadership presentation - blank template or auto-filled from data'
    )
    parser.add_argument('--quarter', type=str, 
                       help='Quarter to analyze (format: YYYY-QN, e.g., 2025-Q4)')
    parser.add_argument('--input', type=str, 
                       help='Path to master_combined.json file')
    parser.add_argument('--output', type=str,
                       help='Output file path (default: templates/Leadership_Presentation_Template.pptx or output/QX_YYYY_Presentation.pptx)')
    parser.add_argument('--blank', action='store_true',
                       help='Force blank template mode even if data provided')
    
    args = parser.parse_args()
    
    # Determine mode: blank template or auto-filled
    data = None
    mode = 'blank'
    
    if args.quarter and args.input and not args.blank:
        # Auto-fill mode
        try:
            # Parse quarter
            year_str, q_str = args.quarter.split('-Q')
            year = int(year_str)
            quarter = int(q_str)
            
            if quarter not in [1, 2, 3, 4]:
                print(f"❌ Error: Quarter must be 1-4, got {quarter}")
                return
            
            # Load data
            data = load_quarterly_data(args.input, year, quarter)
            mode = 'filled'
            
            print(f"\n📊 Loaded {data['quarter_name']} data:")
            print(f"  • Total Issues: {data['total_issues']}")
            print(f"  • Systems: {data['systems_count']}")
            print(f"  • Trend: {data['trend']} ({data['trend_pct']:+.0f}%)")
            
        except Exception as e:
            print(f"❌ Error loading data: {e}")
            print("⚠️  Falling back to blank template mode")
            data = None
            mode = 'blank'
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("\n" + "="*60)
    print(f"Creating Leadership Presentation ({mode.upper()} mode)")
    print("="*60)
    
    print("\n📊 Adding slides:")
    
    add_title_slide(prs, data)
    print("  ✓ Slide 1: Title")
    
    add_executive_summary(prs, data)
    print("  ✓ Slide 2: Executive Summary")
    
    add_metrics_dashboard(prs, data)
    print("  ✓ Slide 3: Key Metrics Dashboard")
    
    add_system_breakdown(prs, data)
    print("  ✓ Slide 4: System Breakdown")
    
    add_success_stories(prs, data)
    print("  ✓ Slide 5: Success Stories")
    
    add_recommendations(prs, data)
    print("  ✓ Slide 6: Recommendations")
    
    add_questions_slide(prs, data)
    print("  ✓ Slide 7: Questions")
    
    # Determine output path
    if args.output:
        output_path = Path(args.output)
    elif data:
        output_path = Path(f'output/{data["quarter_name"].replace(" ", "_")}_Presentation.pptx')
    else:
        output_path = Path('templates/Leadership_Presentation_Template.pptx')
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    
    print("\n" + "="*60)
    print(f"✓ Presentation saved: {output_path}")
    
    if mode == 'blank':
        print("\n📋 Template structure:")
        print("  1. Title slide (customize period and name)")
        print("  2. Executive summary (key achievements)")
        print("  3. Metrics dashboard (6 key metrics)")
        print("  4. System breakdown (chart placeholder)")
        print("  5. Success stories (3 impact examples)")
        print("  6. Recommendations (4 action items)")
        print("  7. Questions (contact info)")
        print("\n🎨 Design: Professional teal and coral palette")
        print("📝 All placeholder text marked with [brackets]")
    else:
        print(f"\n📊 {data['quarter_name']} Presentation:")
        print(f"  • Total Issues: {data['total_issues']}")
        print(f"  • Systems: {data['systems_count']}")
        print(f"  • Trend: {data['trend']} ({data['trend_pct']:+.0f}%)")
        print(f"  • Top System: {data['top_system']} ({data['top_system_count']} issues)")
    
    print("="*60)


if __name__ == '__main__':
    main()

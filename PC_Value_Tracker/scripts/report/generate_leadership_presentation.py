#!/usr/bin/env python3
"""
Generate Leadership Presentation (PowerPoint)

Creates a professional PowerPoint presentation from quarterly data.
Can run in template mode (placeholders) or auto-fill mode (from master.json).

Usage:
    # Auto-fill from data
    python generate_leadership_presentation.py --quarter 2025-Q3

    # Template mode (placeholders only)
    python generate_leadership_presentation.py --template
"""

import json
import sys
from pathlib import Path
from datetime import datetime
from collections import Counter
import argparse

# Fix Windows console encoding
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Color palette - Professional teal and coral
COLORS = {
    'teal': RGBColor(0, 128, 128),
    'teal_dark': RGBColor(0, 102, 102),
    'teal_light': RGBColor(64, 180, 180),
    'coral': RGBColor(255, 127, 80),
    'coral_dark': RGBColor(205, 92, 52),
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
    'gray': RGBColor(128, 128, 128),
    'gray_light': RGBColor(240, 240, 240),
    'green': RGBColor(46, 139, 87),
    'blue': RGBColor(70, 130, 180),
}


def add_title_slide(prs, title, subtitle):
    """Add title slide with teal background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Teal background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['teal']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4), Inches(9), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%B %Y")
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_header(slide, title):
    """Add a section header to a slide."""
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['teal_dark']


def add_metric_card(slide, left, top, width, height, label, value, color=None):
    """Add a metric card with label and value."""
    if color is None:
        color = COLORS['teal']

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    # Value (large)
    value_box = slide.shapes.add_textbox(
        left, top + Inches(0.2), width, Inches(0.8)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Label (smaller)
    label_box = slide.shapes.add_textbox(
        left, top + Inches(0.9), width, Inches(0.5)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER


def add_executive_summary_slide(prs, metrics):
    """Add executive summary slide with key findings."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Executive Summary")

    # Summary text box
    summary_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(5)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True

    # Key achievements
    p = tf.paragraphs[0]
    p.text = "Key Achievements"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['teal_dark']

    achievements = metrics.get('achievements', [
        f"Handled {metrics.get('total_issues', '[X]')} technical support issues",
        f"Resolved {metrics.get('fixed_pct', '[X]')}% of issues (Fixed + Informational)",
        f"Supported {metrics.get('systems_count', '[X]')} different system types",
    ])

    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = f"  * {achievement}"
        p.font.size = Pt(16)
        p.level = 0
        p.space_before = Pt(6)

    # Trends
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Trends & Observations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['teal_dark']
    p.space_before = Pt(18)

    trends = metrics.get('trends', [
        f"Applications stream: {metrics.get('applications_pct', '[X]')}% of workload",
        f"Top system: {metrics.get('top_system', '[System]')} ({metrics.get('top_system_pct', '[X]')}%)",
        f"Compliance-related work: {metrics.get('compliance_pct', '[X]')}% of issues",
    ])

    for trend in trends:
        p = tf.add_paragraph()
        p.text = f"  * {trend}"
        p.font.size = Pt(16)
        p.level = 0
        p.space_before = Pt(6)

    return slide


def add_metrics_dashboard_slide(prs, metrics):
    """Add key metrics dashboard with 6 metric cards."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Key Metrics Dashboard")

    # 6 metric cards in 2 rows of 3
    card_width = Inches(2.8)
    card_height = Inches(1.5)
    start_left = Inches(0.5)
    gap = Inches(0.3)

    row1_top = Inches(1.5)
    row2_top = Inches(3.5)

    cards = [
        ("Total Issues", metrics.get('total_issues', '[X]'), COLORS['teal']),
        ("Fixed/Resolved", f"{metrics.get('fixed_pct', '[X]')}%", COLORS['green']),
        ("Applications", f"{metrics.get('applications_pct', '[X]')}%", COLORS['blue']),
        ("Compliance", f"{metrics.get('compliance_pct', '[X]')}%", COLORS['coral']),
        ("Safety-Related", metrics.get('safety_count', '[X]'), COLORS['coral_dark']),
        ("High Complexity", metrics.get('high_complexity', '[X]'), COLORS['teal_dark']),
    ]

    for i, (label, value, color) in enumerate(cards):
        row = i // 3
        col = i % 3
        left = start_left + col * (card_width + gap)
        top = row1_top if row == 0 else row2_top
        add_metric_card(slide, left, top, card_width, card_height, label, value, color)

    # Footer note
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(9), Inches(0.5)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"Data period: {metrics.get('period', '[Quarter]')}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_system_breakdown_slide(prs, metrics):
    """Add system breakdown slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Work by System")

    # System breakdown table
    systems = metrics.get('systems', [
        ('[System 1]', '[X]', '[X]%'),
        ('[System 2]', '[X]', '[X]%'),
        ('[System 3]', '[X]', '[X]%'),
        ('[System 4]', '[X]', '[X]%'),
        ('[System 5]', '[X]', '[X]%'),
    ])

    # Create visual bars
    start_top = Inches(1.5)
    bar_height = Inches(0.6)
    max_bar_width = Inches(6)
    label_width = Inches(2.5)

    for i, (system, count, pct) in enumerate(systems[:8]):  # Max 8 systems
        top = start_top + i * (bar_height + Inches(0.2))

        # System label
        label_box = slide.shapes.add_textbox(
            Inches(0.5), top, label_width, bar_height
        )
        tf = label_box.text_frame
        tf.paragraphs[0].text = system
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Bar
        try:
            pct_val = float(str(pct).replace('%', '')) / 100
        except:
            pct_val = 0.3

        bar_width = max(Inches(0.3), max_bar_width * pct_val)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.2), top + Inches(0.1),
            bar_width, bar_height - Inches(0.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS['teal']
        bar.line.fill.background()

        # Count/percentage
        value_box = slide.shapes.add_textbox(
            Inches(3.2) + bar_width + Inches(0.1), top,
            Inches(1), bar_height
        )
        tf = value_box.text_frame
        tf.paragraphs[0].text = f"{count} ({pct})"
        tf.paragraphs[0].font.size = Pt(12)

    return slide


def add_stream_breakdown_slide(prs, metrics):
    """Add stream breakdown slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Work by Stream")

    streams = metrics.get('streams', [
        ('Applications', '[X]', '[X]%', 'APO, ACM, DynAMo, Integrity'),
        ('Day-to-Day', '[X]', '[X]%', 'Routine support, troubleshooting'),
        ('Legacy Modernization', '[X]', '[X]%', 'TDC, PLC-5, obsolete systems'),
        ('Diagnostic', '[X]', '[X]%', 'Handed off to other groups'),
        ('After-Hours', '[X]', '[X]%', 'Off-hours support'),
        ('Project', '[X]', '[X]%', 'Capital project handoffs'),
    ])

    start_top = Inches(1.3)
    row_height = Inches(0.85)

    for i, (stream, count, pct, desc) in enumerate(streams):
        top = start_top + i * row_height

        # Stream name and count
        name_box = slide.shapes.add_textbox(
            Inches(0.5), top, Inches(3), Inches(0.4)
        )
        tf = name_box.text_frame
        tf.paragraphs[0].text = f"{stream}: {count} ({pct})"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['teal_dark']

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(0.7), top + Inches(0.35), Inches(8), Inches(0.4)
        )
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = COLORS['gray']

    return slide


def add_highlights_slide(prs, metrics):
    """Add success stories / highlights slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Highlights & Impact")

    highlights = metrics.get('highlights', [
        {
            'title': '[High-Impact Item 1]',
            'system': '[System]',
            'impact': '[Description of value delivered]'
        },
        {
            'title': '[High-Impact Item 2]',
            'system': '[System]',
            'impact': '[Description of value delivered]'
        },
        {
            'title': '[High-Impact Item 3]',
            'system': '[System]',
            'impact': '[Description of value delivered]'
        },
    ])

    start_top = Inches(1.3)
    card_height = Inches(1.5)

    for i, item in enumerate(highlights[:3]):
        top = start_top + i * (card_height + Inches(0.2))

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), top, Inches(9), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['gray_light']
        card.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.7), top + Inches(0.15), Inches(8.5), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = item['title']
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['teal_dark']

        # System tag
        system_box = slide.shapes.add_textbox(
            Inches(0.7), top + Inches(0.5), Inches(2), Inches(0.3)
        )
        tf = system_box.text_frame
        tf.paragraphs[0].text = f"System: {item['system']}"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLORS['coral']

        # Impact
        impact_box = slide.shapes.add_textbox(
            Inches(0.7), top + Inches(0.85), Inches(8.5), Inches(0.5)
        )
        tf = impact_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = item['impact']
        tf.paragraphs[0].font.size = Pt(12)

    return slide


def add_recommendations_slide(prs, metrics):
    """Add recommendations slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    add_section_header(slide, "Recommendations")

    recommendations = metrics.get('recommendations', [
        {
            'title': 'Resource Planning',
            'detail': '[Based on workload data, consider X]'
        },
        {
            'title': 'Training Opportunities',
            'detail': '[Common questions suggest training on X]'
        },
        {
            'title': 'Process Improvements',
            'detail': '[Recurring issues suggest improving X]'
        },
    ])

    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, rec in enumerate(recommendations):
        if i > 0:
            p = tf.add_paragraph()
            p.text = ""

        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{i+1}. {rec['title']}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['teal_dark']
        p.space_before = Pt(12) if i > 0 else Pt(0)

        p = tf.add_paragraph()
        p.text = f"    {rec['detail']}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)

    return slide


def add_questions_slide(prs, contact_name="Process Controls Team", contact_email=""):
    """Add closing questions slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Teal background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['teal']
    background.line.fill.background()

    # Questions?
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(9), Inches(1)
    )
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = contact_name
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if contact_email:
        p = tf.add_paragraph()
        p.text = contact_email
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide


def calculate_metrics(data, quarter):
    """Calculate metrics from master.json data for a specific quarter."""
    # Parse quarter (e.g., "2025-Q3")
    year, q = quarter.split('-Q')
    year = int(year)
    q = int(q)

    # Define quarter months
    quarter_months = {
        1: ['01', '02', '03'],
        2: ['04', '05', '06'],
        3: ['07', '08', '09'],
        4: ['10', '11', '12'],
    }
    months = quarter_months[q]

    # Filter data for the quarter
    quarter_data = []
    for record in data:
        date = record.get('date', '')
        if date and date[:4] == str(year) and date[5:7] in months:
            quarter_data.append(record)

    if not quarter_data:
        return None

    total = len(quarter_data)

    # Resolution breakdown
    resolutions = Counter(r.get('resolution', 'Unknown') for r in quarter_data)
    fixed_count = resolutions.get('Fixed', 0) + resolutions.get('Informational', 0)
    fixed_pct = round(fixed_count / total * 100) if total > 0 else 0

    # Business Impact breakdown
    impacts = Counter(r.get('business_impact', 'Unknown') for r in quarter_data)
    compliance_pct = round(impacts.get('Compliance', 0) / total * 100) if total > 0 else 0
    safety_count = impacts.get('Safety', 0)

    # Stream breakdown
    streams = Counter(r.get('stream', 'Unknown') for r in quarter_data)
    applications_pct = round(streams.get('Applications', 0) / total * 100) if total > 0 else 0

    # System breakdown
    systems = Counter()
    for r in quarter_data:
        sys_name = r.get('system', 'Unknown')
        # Simplify system names (take first part before colon)
        sys_simple = sys_name.split(':')[0].strip() if ':' in sys_name else sys_name
        systems[sys_simple] = systems.get(sys_simple, 0) + 1

    top_systems = systems.most_common(8)
    top_system = top_systems[0][0] if top_systems else 'Unknown'
    top_system_count = top_systems[0][1] if top_systems else 0
    top_system_pct = round(top_system_count / total * 100) if total > 0 else 0

    # Complexity
    complexity = Counter(r.get('complexity', 'Unknown') for r in quarter_data)
    high_complexity = complexity.get('Major', 0) + complexity.get('High', 0)

    # Format streams for slide
    stream_list = []
    stream_descriptions = {
        'Applications': 'APO, ACM, DynAMo, Integrity, Historian',
        'Day-to-Day': 'Routine support, troubleshooting',
        'Legacy Modernization': 'TDC, PLC-5, obsolete systems',
        'Diagnostic': 'Investigated, handed off to other groups',
        'After-Hours': 'Off-hours support, emergencies',
        'Project': 'Capital project handoffs',
    }
    for stream_name, count in streams.most_common():
        pct = round(count / total * 100)
        desc = stream_descriptions.get(stream_name, '')
        stream_list.append((stream_name, count, f"{pct}%", desc))

    # Format systems for slide
    system_list = []
    for sys_name, count in top_systems:
        pct = round(count / total * 100)
        system_list.append((sys_name, count, f"{pct}%"))

    # Find highlights (high complexity or safety items)
    highlights = []
    for r in quarter_data:
        if r.get('complexity') in ['Major', 'High'] or r.get('business_impact') == 'Safety':
            highlights.append({
                'title': r.get('summary', '')[:60] + ('...' if len(r.get('summary', '')) > 60 else ''),
                'system': r.get('system', 'Unknown').split(':')[0].strip(),
                'impact': f"Resolution: {r.get('resolution', 'Unknown')} | Impact: {r.get('business_impact', 'Unknown')}"
            })

    quarter_name = f"Q{q} {year}"
    month_names = {
        1: "January-March",
        2: "April-June",
        3: "July-September",
        4: "October-December"
    }

    return {
        'period': f"{quarter_name} ({month_names[q]} {year})",
        'quarter_name': quarter_name,
        'total_issues': total,
        'fixed_pct': fixed_pct,
        'applications_pct': applications_pct,
        'compliance_pct': compliance_pct,
        'safety_count': safety_count,
        'high_complexity': high_complexity,
        'top_system': top_system,
        'top_system_pct': top_system_pct,
        'systems_count': len(systems),
        'systems': system_list,
        'streams': stream_list,
        'highlights': highlights[:3],
        'achievements': [
            f"Handled {total} technical support issues",
            f"Resolved {fixed_pct}% of issues (Fixed + Informational)",
            f"Supported {len(systems)} different system types",
        ],
        'trends': [
            f"Applications stream: {applications_pct}% of workload",
            f"Top system: {top_system} ({top_system_pct}%)",
            f"Compliance-related work: {compliance_pct}% of issues",
        ],
        'recommendations': [
            {
                'title': 'Resource Planning',
                'detail': f"With {total} issues this quarter, ensure adequate coverage for Applications ({applications_pct}%) and {top_system} support."
            },
            {
                'title': 'Training Opportunities',
                'detail': f"Consider cross-training on {top_system} given its {top_system_pct}% share of workload."
            },
            {
                'title': 'Process Improvements',
                'detail': f"Review {high_complexity} high-complexity items for potential process improvements or automation."
            },
        ],
    }


def generate_presentation(output_path, metrics, title="Process Controls Value Summary"):
    """Generate the full PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build slides
    subtitle = metrics.get('period', '[Quarter/Period]')
    add_title_slide(prs, title, subtitle)
    add_executive_summary_slide(prs, metrics)
    add_metrics_dashboard_slide(prs, metrics)
    add_system_breakdown_slide(prs, metrics)
    add_stream_breakdown_slide(prs, metrics)
    add_highlights_slide(prs, metrics)
    add_recommendations_slide(prs, metrics)
    add_questions_slide(prs, "Process Controls Team")

    # Save
    prs.save(output_path)
    return True


def main():
    parser = argparse.ArgumentParser(
        description='Generate Leadership Presentation (PowerPoint)',
        epilog='Example: python generate_leadership_presentation.py --quarter 2025-Q3'
    )
    parser.add_argument('--quarter', help='Quarter to report on (e.g., 2025-Q3)')
    parser.add_argument('--template', action='store_true', help='Generate template with placeholders only')
    parser.add_argument('--input', default='data/master.json', help='Input JSON file')
    parser.add_argument('--output', help='Output PPTX file (default: auto-generated)')
    parser.add_argument('--verbose', '-v', action='store_true', help='Verbose output')

    args = parser.parse_args()

    if not args.quarter and not args.template:
        print("[ERROR] Must specify --quarter or --template")
        print("  Example: python generate_leadership_presentation.py --quarter 2025-Q3")
        print("  Example: python generate_leadership_presentation.py --template")
        exit(1)

    # Determine output path
    if args.output:
        output_path = Path(args.output)
    elif args.template:
        output_path = Path('output/presentations/Leadership_Presentation_Template.pptx')
    else:
        output_path = Path(f'output/presentations/Leadership_Presentation_{args.quarter}.pptx')

    # Create output directory
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if args.template:
        # Generate template with placeholders
        if args.verbose:
            print("Generating template with placeholders...")
        metrics = {
            'period': '[Quarter/Period]',
            'quarter_name': '[Quarter]',
            'total_issues': '[X]',
            'fixed_pct': '[X]',
            'applications_pct': '[X]',
            'compliance_pct': '[X]',
            'safety_count': '[X]',
            'high_complexity': '[X]',
            'top_system': '[System]',
            'top_system_pct': '[X]',
            'systems_count': '[X]',
        }
    else:
        # Load data and calculate metrics
        if args.verbose:
            print(f"Loading data from {args.input}...")

        input_path = Path(args.input)
        if not input_path.exists():
            print(f"[ERROR] Input file not found: {args.input}")
            exit(1)

        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        if args.verbose:
            print(f"Calculating metrics for {args.quarter}...")

        metrics = calculate_metrics(data, args.quarter)

        if not metrics:
            print(f"[ERROR] No data found for {args.quarter}")
            exit(1)

        if args.verbose:
            print(f"  Found {metrics['total_issues']} issues")

    # Generate presentation
    if args.verbose:
        print("Creating PowerPoint presentation...")

    title = "Process Controls Value Summary"
    generate_presentation(output_path, metrics, title)

    print(f"\n[SUCCESS] Presentation generated!")
    print(f"   {'Template' if args.template else 'Quarter'}: {args.quarter if args.quarter else 'Placeholder'}")
    if not args.template:
        print(f"   Issues: {metrics['total_issues']}")
    print(f"   Output: {output_path}")


if __name__ == '__main__':
    main()

#!/usr/bin/env python3
"""
Add content textboxes to MPC template slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Content for each slide (1-indexed to match slide numbering)
SLIDE_CONTENT = {
    1: {
        "type": "title",  # Title slide - add subtitle and attribution
        "subtitle": "A Case for Readiness",
        "author": "Tony Chiu\nSenior Process Controls Engineer\nJanuary 30, 2026",
        "content": """What I'm Asking For:
* Your assessment of my readiness for Site Lead PCS Engineer role
* Identification of any gaps you see
* Guidance on next steps

What I'm NOT Asking For:
* An immediate promotion decision
* Special treatment or shortcuts

My Goal: Earn your recommendation to move forward"""
    },
    2: {
        "type": "bullets",
        "content": """Based on MPC's expectations, a Site Lead PCS Engineer:

1. Technical Authority: Multi-site platform ownership, escalation point for critical systems
2. Corporate Influence: Shapes enterprise standards, represents site at corporate level
3. Business Leadership: Owns budgets, vendor relationships, strategic planning
4. Crisis Management: Leads high-complexity technical problem-solving
5. Talent Development: Mentors team, builds capability for the future
6. Cross-Functional Leadership: Aligns Operations, Engineering, IT, Safety, Finance

Question for you: Does this align with how you see the role?"""
    },
    3: {
        "type": "bullets",
        "content": """Over the past 18-24 months, my scope has expanded significantly:

* Multi-site technical authority (5 sites: Carson, Wilmington, SLC, Detroit, Anacortes)
* Corporate alarm governance (Tiger Team contributor)
* Platform ownership (Integrity x2, DynAMo x5, Mark VIe)
* Budget/vendor leadership ($2M+ annual, multi-year planning)
* Training programs (Python, SIS, PLC - 200+ hours)
* Crisis leadership (LARINT01 rebuild, Mark VIe restoration)
* Infrastructure modernization (6+ major projects)

I'm not asking to grow into this role-I'm asking for recognition of the work I'm already doing."""
    },
    4: {
        "type": "bullets",
        "content": """My Role:
* Led TWO Integrity platforms (LARRS772 + LARINT01)
* LARINT01: Reference for MPC's 13-refinery fleet
* Rebuilt pipelines, corrected asset mapping, backups

The Impact:
* Carson = enterprise safety model
* 99%+ uptime, audit-ready
* $150K+ cost avoidance

Leadership: Corporate initiative | Multi-site | Standards | Recovery"""
    },
    5: {
        "type": "bullets",
        "content": """My Contributions:
* Active Tiger Team - fixes adopted at multiple sites
* Multi-Site DynAMo SME (5 sites)
* PI-based alarm metrics dashboards
* 30-50% alarm reduction

The Impact:
* Consistent alarm governance across MPC
* Shaped enterprise standards

Leadership: Corporate influence | Multi-site | Standards"""
    },
    6: {
        "type": "bullets",
        "content": """The Challenge:
* GTG chronic unreliability
* Lost Mark VIe expertise

My Role:
* Led GTG troubleshooting, training program, vendor relations

The Impact:
* Zero unplanned trips post-resolution
* Built internal capability

Leadership: Problem solving | Vendor mgmt | Risk mitigation"""
    },
    7: {
        "type": "bullets",
        "content": """A Site Lead is responsible for business and financial leadership.

In 2024-2025 alone, I managed:
* 2025-2029 PCG budget
* Honeywell multi-year license renewals
* Schneider CFA contract
* GE cycle renewals and spares strategy
* Prognost and Trimark service contracts
* Multi-vendor POs, BPOs, and GRs
* PR process optimization with Supply Chain

The Impact:
* Result: Zero licensing lapses, clearer spend forecasting, and stronger vendor accountability
* This financial ownership is a hallmark of the Site Lead PCS role

Leadership Demonstrated:
Business acumen and financial responsibility | Strategic planning (multi-year horizon) | Vendor relationship ownership | Cross-functional coordination (Finance, Supply Chain, Engineering)"""
    },
    8: {
        "type": "bullets",
        "content": """I provide structured development for new and existing team members-another Site Lead PCS expectation.

Key contributions:
* Onboarded new hires (ACM, DynAMo, Integrity, conditional alarming, point deletion)
* Delivered Python training, PLC5/SLC500, SIS courses
* Supported vendor classes and training logistics
* Provided 1:1 mentoring for engineers and technicians
* Created reusable troubleshooting guides and lab environments

The Impact:
* Result: Faster engineer ramp-up, stronger internal talent pipeline, and higher overall controls capability

Leadership Demonstrated:
Investment in team's future capability | Knowledge preservation (prevents tribal knowledge loss) | Strategic thinking (building pipeline, not just solving today's problems)"""
    },
    9: {
        "type": "table",
        "content": """Sites | 5 | Multi-site authority
Systems | 15+ | Platform stewardship
Budget | $2M+ | Financial leadership
Cost Avoid | $150K+ | Business value
Training | 200+ hrs | Talent development
Alarm Cut | 30-50% | Operational improvement
GTG | Zero trips | Risk mitigation
Integrity | 99%+ | Reliability"""
    },
    10: {
        "type": "bullets",
        "content": """Technical:
* DCS, SIS, alarm, APC, OT security
* Crisis leadership, 8+ platforms

Business:
* $2M+ budget, vendor negotiation
* Zero licensing lapses

Leadership:
* Tiger Team, multi-site SME
* Cross-functional (Ops/Safety/IT/Finance)
* Talent development
* HUG presenter, PCE panels

Strategic:
* Multi-year planning (2025-2029)
* 6+ modernization projects"""
    },
    11: {
        "type": "bullets",
        "content": """"Role is new-still defining"
-> I can help shape it based on current work

"Need corporate visibility"
-> Tiger Team + 3 sites + HUG. Where are gaps?

"Not Senior long enough"
-> 18-24 months at Lead level. Tenure target?

"Need PCG alignment"
-> Agreed. Direct or you frame first?

My Ask: Make gaps specific and addressable"""
    },
    12: {
        "type": "bullets",
        "content": """1. Assessment: Where do I stand vs Site Lead?

2. Gaps: Technical? Corporate? Leadership?

3. Path Forward:
   * A: Development plan + timeline
   * B: PCG Technologist discussion
   * C: Recommend to Chris

I Commit To:
* Execute development plan
* Continue Lead-level delivery
* Open to feedback"""
    },
    13: {
        "type": "bullets",
        "content": """Not asking for a new role.
Not asking to grow into it.

Asking for recognition of current leadership.

"Senior PCS Engineer" no longer reflects my work.

Formalizing Site Lead would:
* Enable clearer authority
* Strengthen vendor/cross-site relationships
* Support succession planning
* Align responsibilities with recognition

Thank you. Ready for feedback and next steps."""
    }
}

def add_content_to_presentation(input_file, output_file):
    """Add content textboxes to each slide"""
    prs = Presentation(input_file)
    
    # Slide 0 is the template slide, slides 1-13 are our content slides
    for i in range(1, 14):
        if i >= len(prs.slides):
            print(f"Warning: Slide {i} doesn't exist in presentation")
            continue
            
        slide = prs.slides[i]
        content_data = SLIDE_CONTENT.get(i)
        
        if not content_data:
            print(f"No content defined for slide {i}")
            continue
        
        # Add content based on type
        if content_data.get("type") == "title":
            # Title slide - add subtitle and author
            # Subtitle below title
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content_data["subtitle"]
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.bold = True
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            # Author info
            author_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.6)
            )
            author_frame = author_box.text_frame
            author_frame.text = content_data["author"]
            author_para = author_frame.paragraphs[0]
            author_para.font.size = Pt(16)
            author_para.alignment = PP_ALIGN.CENTER
            
            # Content bullets below author
            if "content" in content_data:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(3.3), Inches(9.0), Inches(2.7)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                lines = content_data["content"].strip().split('\n')
                for j, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        continue
                    
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    is_bullet = line.startswith('*')
                    if is_bullet:
                        p.text = line[1:].strip()
                        p.level = 0
                        p.font.size = Pt(14)
                    else:
                        p.text = line
                        p.level = 0
                        if line.endswith(':'):
                            p.font.bold = True
                            p.font.size = Pt(16)
                        else:
                            p.font.size = Pt(14)
                    p.space_after = Pt(6)
            
        elif content_data.get("type") == "table":
            # Create table for slide 10
            rows = content_data["content"].strip().split('\n')
            content_box = slide.shapes.add_textbox(
                Inches(0.42), Inches(1.3), Inches(9.15), Inches(4.3)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for row in rows:
                parts = [p.strip() for p in row.split('|')]
                if len(parts) == 3:
                    p = text_frame.add_paragraph()
                    p.text = f"{parts[0]}: {parts[1]}"
                    p.font.size = Pt(13)
                    p.space_after = Pt(8)
                    p.level = 0
                    
                    # Add significance as indented text
                    p2 = text_frame.add_paragraph()
                    p2.text = f"({parts[2]})"
                    p2.font.size = Pt(11)
                    p2.font.italic = True
                    p2.space_after = Pt(12)
                    p2.level = 1
            
        else:
            # Regular bullet content
            content_box = slide.shapes.add_textbox(
                Inches(0.42), Inches(1.3), Inches(9.15), Inches(4.3)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            # Split content into lines and add paragraphs
            lines = content_data["content"].strip().split('\n')
            for j, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                if j == 0:
                    # First paragraph uses existing text frame
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Determine if it's a bullet point
                is_bullet = line.startswith('*') or line.startswith('-')
                is_numbered = line and line[0].isdigit() and '. ' in line[:4]
                
                if is_bullet:
                    p.text = line[1:].strip()  # Remove bullet character
                    p.level = 0
                    p.font.size = Pt(14)
                elif is_numbered:
                    p.text = line
                    p.level = 0
                    p.font.size = Pt(14)
                else:
                    # Section header or regular text
                    p.text = line
                    p.level = 0
                    if line.endswith(':') or line.isupper():
                        p.font.bold = True
                        p.font.size = Pt(16)
                        p.space_after = Pt(6)
                    else:
                        p.font.size = Pt(14)
                
                p.space_after = Pt(8)
        
        print(f"* Slide {i}: {content_data.get('type', 'bullets')}")
    
    prs.save(output_file)
    print(f"\nPresentation saved to: {output_file}")

if __name__ == "__main__":
    import os
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    presentations_dir = os.path.join(parent_dir, "Presentations")
    
    add_content_to_presentation(
        os.path.join(presentations_dir, "titles-only.pptx"),
        os.path.join(presentations_dir, "Site_Lead_PCS_Engineer_Presentation_01_30_26.pptx")
    )

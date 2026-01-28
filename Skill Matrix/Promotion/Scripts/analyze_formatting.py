"""Analyze presentation formatting against established standards."""
from pptx import Presentation
from pptx.util import Pt

prs = Presentation('Presentations/Site_Lead_PCS_Engineer_Presentation_01_30_26.pptx')

# Formatting standards
STANDARDS = {
    'level_0': 16,  # Level 0 bullets
    'level_1': 15,  # Level 1 bullets
    'level_2': 14,  # Level 2+ bullets
    'line_spacing': 1.25
}

non_compliant = []
total_paragraphs = 0

for slide_idx, slide in enumerate(prs.slides):
    slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx}"
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    total_paragraphs += 1
                    
                    for run in paragraph.runs:
                        if run.font.size:
                            actual_size = round(run.font.size.pt, 1)
                            level = paragraph.level
                            
                            # Determine expected size
                            if level == 0:
                                expected = STANDARDS['level_0']
                            elif level == 1:
                                expected = STANDARDS['level_1']
                            else:
                                expected = STANDARDS['level_2']
                            
                            # Allow 0.5pt tolerance for rounding
                            if abs(actual_size - expected) > 0.5:
                                non_compliant.append({
                                    'slide': slide_idx,
                                    'title': slide_title,
                                    'level': level,
                                    'actual': actual_size,
                                    'expected': expected,
                                    'text': run.text[:60].strip()
                                })

print(f"Formatting Analysis Report")
print(f"=" * 80)
print(f"Total slides analyzed: {len(prs.slides)}")
print(f"Total text paragraphs: {total_paragraphs}")
print(f"Formatting deviations found: {len(non_compliant)}")
print()

if non_compliant:
    print("Deviations from formatting standards:")
    print("-" * 80)
    
    # Group by slide
    by_slide = {}
    for item in non_compliant:
        slide_num = item['slide']
        if slide_num not in by_slide:
            by_slide[slide_num] = []
        by_slide[slide_num].append(item)
    
    for slide_num in sorted(by_slide.keys()):
        items = by_slide[slide_num]
        print(f"\nSlide {slide_num}: {items[0]['title']}")
        for item in items[:5]:  # Show first 5 per slide
            print(f"  • Level {item['level']}: {item['actual']}pt (expected {item['expected']}pt)")
            print(f"    Text: '{item['text']}'")
        if len(items) > 5:
            print(f"  ... and {len(items) - 5} more deviations on this slide")
    
    print()
    print(f"Summary: {len(non_compliant)} total deviations across {len(by_slide)} slides")
else:
    print("✓ All text formatting complies with standards!")
    print(f"  - Level 0 bullets: {STANDARDS['level_0']}pt")
    print(f"  - Level 1 bullets: {STANDARDS['level_1']}pt")
    print(f"  - Level 2+ bullets: {STANDARDS['level_2']}pt")

#!/usr/bin/env python3
"""
Standardize body text formatting in PowerPoint presentation
Makes all body text consistent, professional, and readable
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def standardize_body_text_formatting(input_file, output_file):
    """
    Apply consistent formatting to all body text shapes (non-title content)
    
    Standard formatting:
    - Font size: 14pt for body text, 12pt for sub-bullets
    - Font: Calibri (or template default)
    - Line spacing: 1.15
    - Space after paragraph: 6pt
    - Bullets: Proper indentation levels
    - Alignment: Left for bullets, preserved for others
    """
    prs = Presentation(input_file)
    
    changes_made = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Skip first slide (title slide - slide 0 in 0-indexed)
        if slide_num == 1:
            continue
            
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            
            # Identify if this is likely a title shape (top of slide, larger font)
            is_likely_title = (
                shape.top < 1000000 and  # Top inch of slide (in EMUs)
                len(text_frame.paragraphs) <= 2 and
                any(p.font.size and p.font.size > Pt(20) for p in text_frame.paragraphs if p.font.size)
            )
            
            # Skip title shapes
            if is_likely_title:
                continue
            
            # Apply standard formatting to body text
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                # Skip empty paragraphs
                if not paragraph.text.strip():
                    continue
                
                # Determine if this is a section header within body text
                is_section_header = (
                    paragraph.text.endswith(':') or
                    paragraph.font.bold or
                    (para_idx == 0 and not paragraph.level)
                )
                
                # Standard font size based on content type (increased for readability)
                if is_section_header:
                    paragraph.font.size = Pt(18)
                    paragraph.font.bold = True
                    paragraph.space_after = Pt(8)
                elif paragraph.level == 0:
                    paragraph.font.size = Pt(16)
                    paragraph.space_after = Pt(8)
                elif paragraph.level == 1:
                    paragraph.font.size = Pt(15)
                    paragraph.space_after = Pt(6)
                else:
                    paragraph.font.size = Pt(14)
                    paragraph.space_after = Pt(6)
                
                # Line spacing (increased for readability)
                paragraph.line_spacing = 1.25
                
                # Alignment: left for bullets, preserve for non-bullets
                if paragraph.level is not None and paragraph.level >= 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                
                # Space before (only for section headers)
                if is_section_header and para_idx > 0:
                    paragraph.space_before = Pt(12)
                else:
                    paragraph.space_before = Pt(0)
                
                changes_made += 1
        
        print(f"Slide {slide_num}: Formatted")
    
    prs.save(output_file)
    print(f"\nFormatting complete!")
    print(f"Total paragraphs formatted: {changes_made}")
    print(f"Output saved to: {output_file}")

if __name__ == "__main__":
    import os
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    presentations_dir = os.path.join(parent_dir, "Presentations")
    pptx_file = os.path.join(presentations_dir, "Site_Lead_PCS_Engineer_Presentation_01_30_26.pptx")
    
    standardize_body_text_formatting(pptx_file, pptx_file)
